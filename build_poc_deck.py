"""
Build the Hyreo internal deck:
'Building Signature Sprint with Claude.ai — Effort, Scale, and a Path Forward'
Audience: internal Hyreo team. Compares BMAD vs CCPM and recommends CCPM.
Effort numbers are defensible estimates (clearly labeled).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x2A, 0x3D)
NAVY_2  = RGBColor(0x2A, 0x3A, 0x52)
CYAN    = RGBColor(0x3B, 0xD4, 0xE7)
CYAN_2  = RGBColor(0x0E, 0xA5, 0xB8)
INK     = RGBColor(0x1F, 0x2A, 0x37)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
SOFT    = RGBColor(0xF5, 0xF9, 0xFC)
ACCENT  = RGBColor(0xE8, 0xFA, 0xFD)
RULE    = RGBColor(0xE5, 0xEA, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
AMBER   = RGBColor(0xF5, 0xA6, 0x23)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
RED     = RGBColor(0xE7, 0x4C, 0x3C)

FONT     = "Calibri"
DISPLAY  = "Calibri"
MONO     = "Consolas"

# ── Setup: 16:9 ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(BLANK)

def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.text_frame.margin_left = shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = shape.text_frame.margin_bottom = 0
    return shape

def add_round_rect(slide, x, y, w, h, fill, line=None, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.text_frame.margin_left = Inches(0.12)
    shape.text_frame.margin_right = Inches(0.12)
    shape.text_frame.margin_top = Inches(0.06)
    shape.text_frame.margin_bottom = Inches(0.06)
    return shape

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=INK, font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs: list of (text, dict). dict keys: size, bold, italic, color, font."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for i, (text, opts) in enumerate(runs):
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", FONT)
        r.font.size = Pt(opts.get("size", 14))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", INK)
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=INK, gap=8, lead_color=NAVY, line_spacing=1.25):
    """items: list of strings OR list of (lead, rest) tuples for bold lead text."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(gap)
        # bullet glyph
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name = FONT
        rb.font.size = Pt(size)
        rb.font.color.rgb = CYAN
        rb.font.bold = True
        if isinstance(item, tuple):
            lead, rest = item
            rl = p.add_run(); rl.text = lead
            rl.font.name = FONT; rl.font.size = Pt(size); rl.font.bold = True; rl.font.color.rgb = lead_color
            rr = p.add_run(); rr.text = rest
            rr.font.name = FONT; rr.font.size = Pt(size); rr.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
    return tb

def add_hyreo(slide, x, y, size=44):
    """Render the 'hyreo' wordmark: navy 'hy' + cyan 'reo'."""
    w = Inches(2.4); h = Inches(0.9)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "hy"
    r1.font.name = DISPLAY; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = "reo"
    r2.font.name = DISPLAY; r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = CYAN
    return tb

def add_page_chrome(slide, eyebrow, title, page_no=None, total=None):
    # top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.18), NAVY)
    # tiny cyan stripe
    add_rect(slide, 0, Inches(0.18), Inches(2.0), Inches(0.05), CYAN)
    # mini wordmark, top-right
    tb = slide.shapes.add_textbox(SW - Inches(1.5), Inches(0.32), Inches(1.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run(); r1.text = "hy"; r1.font.name = DISPLAY; r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = "reo"; r2.font.name = DISPLAY; r2.font.size = Pt(16); r2.font.bold = True; r2.font.color.rgb = CYAN
    # eyebrow
    add_text(slide, Inches(0.55), Inches(0.32), Inches(8), Inches(0.3),
             eyebrow.upper(), size=10, bold=True, color=CYAN)
    # title
    add_text(slide, Inches(0.55), Inches(0.62), Inches(11.5), Inches(0.7),
             title, size=26, bold=True, color=NAVY)
    # divider rule
    add_rect(slide, Inches(0.55), Inches(1.32), Inches(12.2), Emu(9525), RULE)
    # page number
    if page_no is not None and total is not None:
        add_text(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.3),
                 f"{page_no:02d} / {total:02d}",
                 size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def add_footer_brand(slide):
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             "Hyreo Labs  ·  Internal  ·  Building with Claude.ai",
             size=9, italic=True, color=MUTED)

# ── Table helper ─────────────────────────────────────────────────────────────
def add_clean_table(slide, x, y, w, h, rows, col_widths=None, header=True,
                    header_fill=NAVY, header_text=WHITE, alt_fill=None,
                    body_size=11, header_size=11):
    n_rows = len(rows); n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT
            is_header = header and r_idx == 0
            r.font.size = Pt(header_size if is_header else body_size)
            r.font.bold = is_header
            r.font.color.rgb = header_text if is_header else INK
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = header_fill
            else:
                if alt_fill is not None and r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.fore_color.rgb = WHITE
    return tbl_shape

# =============================================================================
# Slides
# =============================================================================
TOTAL = 15

# ── 1. Cover ─────────────────────────────────────────────────────────────────
s = add_slide()
add_rect(s, 0, 0, SW, SH, WHITE)
# left navy panel
add_rect(s, 0, 0, Inches(4.6), SH, NAVY)
# cyan accent stripe
add_rect(s, Inches(4.6), 0, Inches(0.08), SH, CYAN)
# wordmark on left panel (white + cyan)
tb = s.shapes.add_textbox(Inches(0.55), Inches(0.6), Inches(3), Inches(1.0))
p = tb.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = "hy"; r1.font.name = DISPLAY; r1.font.size = Pt(48); r1.font.bold = True; r1.font.color.rgb = WHITE
r2 = p.add_run(); r2.text = "reo"; r2.font.name = DISPLAY; r2.font.size = Pt(48); r2.font.bold = True; r2.font.color.rgb = CYAN
add_text(s, Inches(0.55), Inches(1.6), Inches(4), Inches(0.3),
         "Hyreo Labs", size=12, color=CYAN)
# vertical decorative dots on left
for i in range(3):
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55 + i * 0.25), Inches(6.6), Inches(0.12), Inches(0.12))
    dot.shadow.inherit = False
    dot.fill.solid(); dot.fill.fore_color.rgb = CYAN if i == 0 else (RGBColor(0x4B, 0x5E, 0x7A) if i == 1 else RGBColor(0x33, 0x45, 0x60))
    dot.line.fill.background()

add_text(s, Inches(0.55), Inches(6.9), Inches(4), Inches(0.3),
         "Internal  ·  Confidential", size=9, italic=True, color=RGBColor(0x9F, 0xB2, 0xC9))

# right side title block
add_text(s, Inches(5.0), Inches(1.5), Inches(8), Inches(0.4),
         "INTERNAL TEAM PRESENTATION  ·  APRIL 2026",
         size=11, bold=True, color=CYAN)
add_text(s, Inches(5.0), Inches(2.05), Inches(8), Inches(1.6),
         "Building with Claude.ai", size=44, bold=True, color=NAVY)
add_text(s, Inches(5.0), Inches(3.15), Inches(8), Inches(1.0),
         "Effort, scalability, and a path forward for the\nSignature Sprint POC and beyond.",
         size=18, italic=True, color=MUTED)

# meta block
add_rect(s, Inches(5.0), Inches(4.7), Inches(7.6), Emu(9525), RULE)
meta = [("PROJECT",    "Signature Sprint  ·  v2.7  ·  POC"),
        ("PREPARED BY","Hyreo Labs"),
        ("AUDIENCE",   "Hyreo internal team"),
        ("FOCUS",      "Effort comparison · BMAD vs CCPM · Onboarding")]
for i, (k, v) in enumerate(meta):
    y = Inches(4.85 + i * 0.42)
    add_text(s, Inches(5.0), y, Inches(2.0), Inches(0.3), k, size=9.5, bold=True, color=MUTED)
    add_text(s, Inches(7.0), y, Inches(6.0), Inches(0.3), v, size=12, color=INK)

add_text(s, Inches(5.0), Inches(7.05), Inches(7), Inches(0.3),
         "Hyreo Labs  ·  hyreo.com", size=9, italic=True, color=MUTED)

# ── 2. Agenda ────────────────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "Agenda", "What we'll cover", page_no=2, total=TOTAL)
agenda = [
    ("01.  The challenge",                "Where Signature Sprint started and the constraints we worked under."),
    ("02.  How Claude.ai helped",         "The role it played across planning, building, and writing."),
    ("03.  Effort comparison",            "Defensible estimates: with AI vs without AI, by phase."),
    ("04.  What scaled — and what didn't","Honest read on where the leverage came from."),
    ("05.  Applicability to the project", "Concrete next features to ship with the same workflow."),
    ("06.  Method: BMAD vs CCPM",         "Side-by-side comparison."),
    ("07.  Recommendation",               "Why CCPM fits Signature Sprint right now."),
    ("08.  Onboarding a new developer",   "Day 1 / Week 1 / Month 1 with CCPM."),
    ("09.  Risks, guardrails, roadmap",   "What to watch for; what comes next."),
]
y = Inches(1.6)
for label, sub in agenda:
    add_text(s, Inches(0.6), y, Inches(4.5), Inches(0.4),
             label, size=14, bold=True, color=NAVY)
    add_text(s, Inches(5.2), y, Inches(7.6), Inches(0.4),
             sub, size=12, color=MUTED)
    y += Inches(0.52)
add_footer_brand(s)

# ── 3. The challenge ─────────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "01  ·  Context", "The challenge we set ourselves", page_no=3, total=TOTAL)
add_runs(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.2),
         [("Build an end-to-end behavioral assessment POC ", {"size": 16, "color": INK}),
          ("— two themed UIs, a four-level scoring engine, a manager dashboard, and a full ", {"size": 16, "color": INK}),
          ("anti-faking", {"size": 16, "color": NAVY, "bold": True}),
          (" validation layer — fast enough to put in front of evaluators in weeks, not quarters.", {"size": 16, "color": INK})],
         line_spacing=1.3)

# Three constraint cards
cards = [
    ("Time pressure",  "Single sprint to demo; multi-week budget, not multi-quarter."),
    ("Small team",     "Effectively one engineer carrying frontend, scoring, and docs."),
    ("Quality bar",    "Has to look and feel like a Hyreo product, not a prototype."),
]
cw = Inches(3.95); gap = Inches(0.18); x0 = Inches(0.6); y0 = Inches(3.4)
for i, (t, sub) in enumerate(cards):
    x = x0 + (cw + gap) * i
    add_round_rect(s, x, y0, cw, Inches(2.0), SOFT, line=RULE)
    # cyan tag
    add_rect(s, x + Inches(0.25), y0 + Inches(0.25), Inches(0.18), Inches(0.18), CYAN)
    add_text(s, x + Inches(0.55), y0 + Inches(0.18), cw - Inches(0.7), Inches(0.4),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y0 + Inches(0.85), cw - Inches(0.5), Inches(1.0),
             sub, size=12, color=INK)

# Outcome strip
add_round_rect(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(1.0), ACCENT, line=CYAN, radius=0.18)
add_runs(s, Inches(0.85), Inches(5.85), Inches(11.8), Inches(0.7),
         [("Outcome.  ", {"size": 14, "bold": True, "color": NAVY}),
          ("v2.7 ships with both themes, a working scoring engine with validation flags, full PRD/User-Guide/Report-Reference docs, and a print-ready manager report — running in any modern browser without a backend.",
           {"size": 13, "color": INK})],
         line_spacing=1.25)
add_footer_brand(s)

# ── 4. How Claude.ai helped ──────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "02  ·  Approach", "Claude.ai as a development partner", page_no=4, total=TOTAL)
add_text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
         "Five concrete roles Claude played across the build.",
         size=13, italic=True, color=MUTED)

roles = [
    ("Architect",  "Mapped competency framework → dimensions → BARS → scoring formula in one pass."),
    ("Pair coder", "Wrote the scoring engine, both UI themes, and the manager dashboard alongside us."),
    ("Refactor",   "Renamed Catalyst Protocol → Signature Sprint cleanly without breaking internals."),
    ("Doc writer", "Produced PRODUCT-DOC, USER-GUIDE, REPORT-REFERENCE in a consistent voice."),
    ("Reviewer",   "Caught regressions when L2 rules and skipped-question logic changed in v2.7."),
]
# Two-column grid
x_left = Inches(0.6); x_right = Inches(6.95)
card_w = Inches(6.0); card_h = Inches(1.55)
positions = [
    (x_left,  Inches(2.1)),
    (x_right, Inches(2.1)),
    (x_left,  Inches(3.8)),
    (x_right, Inches(3.8)),
    (x_left,  Inches(5.5)),
]
for (x, y), (title, body) in zip(positions, roles):
    add_round_rect(s, x, y, card_w, card_h, WHITE, line=RULE, radius=0.1)
    add_rect(s, x, y, Inches(0.12), card_h, CYAN)
    add_text(s, x + Inches(0.3), y + Inches(0.18), card_w - Inches(0.5), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(0.95),
             body, size=12, color=INK)

# Right side last-card placeholder text
add_round_rect(s, x_right, Inches(5.5), card_w, card_h, NAVY, line=NAVY, radius=0.1)
add_text(s, x_right + Inches(0.3), Inches(5.65), card_w - Inches(0.5), Inches(0.4),
         "The shift", size=15, bold=True, color=CYAN)
add_text(s, x_right + Inches(0.3), Inches(6.05), card_w - Inches(0.5), Inches(1.0),
         "Claude moved us from \"engineer drafts everything\" to \"engineer reviews, decides, and merges.\"",
         size=12, italic=True, color=WHITE)
add_footer_brand(s)

# ── 5. Effort comparison ─────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "03  ·  Effort", "With AI vs without AI — by phase",
                page_no=5, total=TOTAL)
# Caption / disclaimer
add_runs(s, Inches(0.6), Inches(1.5), Inches(12.2), Inches(0.4),
         [("Numbers below are ", {"size": 11, "italic": True, "color": MUTED}),
          ("defensible estimates", {"size": 11, "italic": True, "bold": True, "color": NAVY}),
          (" (calibrated against this build, not stopwatch-measured). Use as order-of-magnitude, not exact.",
           {"size": 11, "italic": True, "color": MUTED})])

rows = [
    ["Phase",                              "Without AI",       "With Claude.ai", "Speed-up"],
    ["Discovery & competency framework",   "3–4 days",         "0.5 day",         "≈ 6×"],
    ["Scenario design (4 levels, ~21 items)","5–7 days",       "1.5 days",        "≈ 4×"],
    ["Scoring engine + BARS mapping",      "4–5 days",         "1 day",           "≈ 4×"],
    ["Adventure UI (Theme 1)",             "5–6 days",         "1.5 days",        "≈ 4×"],
    ["Mission Control UI (Theme 2)",       "4–5 days",         "1 day",           "≈ 4×"],
    ["Manager dashboard + report",         "3–4 days",         "1 day",           "≈ 3–4×"],
    ["Anti-faking validation flags",       "3 days",           "0.5–1 day",       "≈ 3×"],
    ["Brand rename (Catalyst → Signature)","1.5 days",         "≈ 2 hrs",         "≈ 6×"],
    ["Documentation (PRD + guides + report ref)","4–5 days",   "1 day",           "≈ 4–5×"],
    ["Total (effort, not calendar)",       "≈ 33–43 days",     "≈ 8–9 days",      "≈ 4–5×"],
]
table = add_clean_table(
    s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(4.6), rows,
    col_widths=[Inches(4.8), Inches(2.5), Inches(2.7), Inches(2.2)],
    alt_fill=SOFT, body_size=11, header_size=11.5,
)
# bold the totals row visually by re-styling its first cell with NAVY-on-cyan tone
tbl = table.table
last = tbl.rows[len(rows) - 1]
for i, cell in enumerate(last.cells):
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = NAVY

# Take-aways below
add_round_rect(s, Inches(0.6), Inches(6.6), Inches(12.2), Inches(0.6), NAVY, line=NAVY, radius=0.18)
add_runs(s, Inches(0.85), Inches(6.7), Inches(11.8), Inches(0.5),
         [("Headline.  ", {"size": 13, "bold": True, "color": CYAN}),
          ("A ~5–6 week build compresses to ~1.5–2 weeks of focused engineering effort. Calendar time depends on review and stakeholder cycles, not typing.",
           {"size": 12, "color": WHITE})],
         anchor=MSO_ANCHOR.MIDDLE)

# ── 6. What scaled — and what didn't ─────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "04  ·  Honest read", "What scaled — and what didn't",
                page_no=6, total=TOTAL)

# Two columns
col_w = Inches(6.0); gap = Inches(0.2)
# Left — scaled well (green accent)
add_round_rect(s, Inches(0.6), Inches(1.6), col_w, Inches(5.4), WHITE, line=RULE, radius=0.06)
add_rect(s, Inches(0.6), Inches(1.6), col_w, Inches(0.5), GREEN)
add_text(s, Inches(0.8), Inches(1.65), col_w - Inches(0.4), Inches(0.4),
         "What scaled", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
left_items = [
    ("Repeatable prompts.",  " Scoring rules, BARS, dimension weights — encoded once, reused everywhere."),
    ("Shared context files.","  PRODUCT-DOC, USER-GUIDE, REPORT-REFERENCE feed every new task instantly."),
    ("Parallel UI work.",    " Adventure and Mission built against the same shared engine without diverging."),
    ("Refactor under cover.","  The Catalyst → Signature Sprint rename happened in hours, not days."),
    ("Documentation parity.","  Code and docs stay aligned because both come out of the same workflow."),
]
add_bullets(s, Inches(0.85), Inches(2.25), col_w - Inches(0.45), Inches(4.7),
            left_items, size=12, gap=8, lead_color=NAVY, line_spacing=1.25)

# Right — limits (amber accent)
x_r = Inches(0.6) + col_w + gap
add_round_rect(s, x_r, Inches(1.6), col_w, Inches(5.4), WHITE, line=RULE, radius=0.06)
add_rect(s, x_r, Inches(1.6), col_w, Inches(0.5), AMBER)
add_text(s, x_r + Inches(0.2), Inches(1.65), col_w - Inches(0.4), Inches(0.4),
         "What didn't (yet)", size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
right_items = [
    ("Subjective design.",   " Visual taste and animation feel still need a human in the loop."),
    ("Behavioral validity.", " Statistical rigor of validation flags needs an I/O psych review."),
    ("Edge-case bugs.",      " Timer / skip / state-restore corners required hands-on debugging."),
    ("Code review still matters.", " AI moves fast; merging without review compounds risk fast."),
    ("Onboarding artifacts.","  Without CCPM-style structure, context is in our heads, not the repo."),
]
add_bullets(s, x_r + Inches(0.25), Inches(2.25), col_w - Inches(0.45), Inches(4.7),
            right_items, size=12, gap=8, lead_color=NAVY, line_spacing=1.25)
add_footer_brand(s)

# ── 7. Applicability to the project ──────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "05  ·  Applicability", "Where the same workflow earns its keep next",
                page_no=7, total=TOTAL)

add_text(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.4),
         "Concrete next features for Signature Sprint that fit the same AI-paired workflow.",
         size=12, italic=True, color=MUTED)

apps = [
    ("New competencies",     "Add Collaboration, Customer Focus, Learning Agility — same pattern: dimensions → BARS → scenarios → scoring weights."),
    ("Persistence layer",    "Swap browser localStorage for a real backend (Postgres + REST). Engine surface stays unchanged."),
    ("Multi-evaluator review","Add roles, comments, and a shortlist queue on top of the manager portal."),
    ("Report exports",       "PDF / CSV bulk export from the candidate list; structured JSON for ATS handoff."),
    ("Localization",         "Hindi / regional language scenario sets; one prompt, one review pass per locale."),
    ("Analytics dashboard",  "Cohort views — score distribution, flag rates, dimension means by role band."),
]
# 3x2 grid
cw = Inches(4.0); ch = Inches(1.65); gap_x = Inches(0.15); gap_y = Inches(0.18)
x0 = Inches(0.6); y0 = Inches(2.0)
for i, (t, body) in enumerate(apps):
    col = i % 3; row = i // 3
    x = x0 + col * (cw + gap_x); y = y0 + row * (ch + gap_y)
    add_round_rect(s, x, y, cw, ch, WHITE, line=RULE, radius=0.08)
    add_rect(s, x, y, Inches(0.1), ch, CYAN)
    add_text(s, x + Inches(0.25), y + Inches(0.15), cw - Inches(0.4), Inches(0.4),
             t, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.55), cw - Inches(0.4), Inches(1.05),
             body, size=11.5, color=INK)

add_round_rect(s, Inches(0.6), Inches(5.7), Inches(12.2), Inches(1.05), SOFT, line=RULE, radius=0.1)
add_runs(s, Inches(0.85), Inches(5.85), Inches(11.8), Inches(0.85),
         [("The pattern.  ", {"size": 13, "bold": True, "color": NAVY}),
          ("Each item above is a vertical slice — spec, scenarios, code, and docs — that fits a single CCPM epic. Sized for one engineer + Claude over 3–8 days, end to end.",
           {"size": 12, "color": INK})],
         line_spacing=1.3)
add_footer_brand(s)

# ── 8. BMAD vs CCPM side-by-side ─────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "06  ·  Method", "BMAD vs CCPM — side by side",
                page_no=8, total=TOTAL)

add_text(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.4),
         "Two mature AI-development methods. Same goal — repeatable, structured AI-paired delivery — different shapes.",
         size=12, italic=True, color=MUTED)

rows = [
    ["Dimension",         "BMAD-METHOD",                                            "CCPM"],
    ["Core idea",         "Specialized agent personas drive a phased pipeline",     "GitHub Issues drive parallel sub-agent work"],
    ["Roles",             "Analyst → PM → Architect → SM → Dev → QA",               "PRD → Epic → Tasks → Issues → Implementation"],
    ["Unit of work",      "Story file with full context for the Dev agent",          "GitHub Issue worked in its own git worktree"],
    ["Parallelism",       "Sequential by design; one phase at a time",              "Parallel by design; many issues / agents at once"],
    ["Best fit",          "Greenfield — building a new product from scratch",       "Brownfield — scaling and shipping an existing product"],
    ["Team coordination", "Single dev + agents; story files are the handoff",       "Team-first; GitHub is the source of truth"],
    ["Onboarding cost",   "Higher — new vocabulary + persona files",                "Lower — anyone who knows GitHub Issues can join"],
    ["Failure mode",      "Heavy ceremony if the project is small",                 "Issue sprawl without good epic hygiene"],
]
table = add_clean_table(
    s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(4.8), rows,
    col_widths=[Inches(2.6), Inches(4.8), Inches(4.8)],
    alt_fill=SOFT, body_size=10.5, header_size=11.5,
)
# Lightly tint method columns
tbl = table.table
for r_idx in range(1, len(rows)):
    if r_idx % 2 == 1:  # already alt-filled SOFT, skip
        continue

# Verdict strip
add_round_rect(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.5), NAVY, line=NAVY, radius=0.2)
add_runs(s, Inches(0.85), Inches(6.9), Inches(11.8), Inches(0.4),
         [("Read.  ", {"size": 12, "bold": True, "color": CYAN}),
          ("BMAD optimizes for ", {"size": 12, "color": WHITE}),
          ("starting", {"size": 12, "color": WHITE, "italic": True}),
          (".  CCPM optimizes for ", {"size": 12, "color": WHITE}),
          ("scaling", {"size": 12, "color": WHITE, "italic": True}),
          (".  Signature Sprint is past start.", {"size": 12, "color": WHITE})],
         anchor=MSO_ANCHOR.MIDDLE)

# ── 9. Recommendation ────────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "07  ·  Recommendation", "We recommend CCPM for Signature Sprint",
                page_no=9, total=TOTAL)

# Big stamp
add_round_rect(s, Inches(0.6), Inches(1.6), Inches(4.5), Inches(2.4), NAVY, line=NAVY, radius=0.06)
add_text(s, Inches(0.85), Inches(1.75), Inches(4.0), Inches(0.4),
         "VERDICT", size=11, bold=True, color=CYAN)
add_text(s, Inches(0.85), Inches(2.1), Inches(4.0), Inches(0.8),
         "Adopt CCPM", size=34, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(3.0), Inches(4.0), Inches(1.0),
         "Lightweight overlay on top of GitHub.\nNo new vocabulary. Parallel by design.",
         size=12, italic=True, color=RGBColor(0xC7, 0xD7, 0xEA))

# Why CCPM (right column)
why_x = Inches(5.4); why_w = Inches(7.4)
add_text(s, why_x, Inches(1.6), why_w, Inches(0.4),
         "Why CCPM, specifically", size=15, bold=True, color=NAVY)
why = [
    ("Project is past greenfield.",      " The system exists; the work is iterate-and-scale, where CCPM shines."),
    ("Team is small and growing.",       " GitHub Issues are an onboarding artifact — pick one, read it, ship it."),
    ("Parallel work fits naturally.",    " Each next feature (persistence, locales, analytics) is an independent vertical."),
    ("Single source of truth.",          " PRDs and epics live next to the code, not in slides or chat threads."),
    ("Low ceremony, high structure.",    " Dev keeps current tools; AI plugs in via worktrees and sub-agents."),
]
add_bullets(s, why_x, Inches(2.05), why_w, Inches(4.5), why,
            size=12, gap=10, lead_color=NAVY, line_spacing=1.25)

# Caveat strip
add_round_rect(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.7), ACCENT, line=CYAN, radius=0.18)
add_runs(s, Inches(0.85), Inches(6.55), Inches(11.8), Inches(0.6),
         [("When BMAD wins instead.  ", {"size": 12, "bold": True, "color": NAVY}),
          ("If we kick off a brand-new product (e.g., a separate Hyreo Compass or a 360° tool) from a clean slate, BMAD's analyst → architect → dev pipeline is a stronger starting frame.",
           {"size": 12, "color": INK})],
         anchor=MSO_ANCHOR.MIDDLE)

# ── 10. CCPM workflow ────────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "08  ·  CCPM workflow", "How a feature flows under CCPM",
                page_no=10, total=TOTAL)

# 5 stage flow
stages = [
    ("PRD",        "Product brief\nin /prds/"),
    ("Epic",       "Decomposed into\nphased plan"),
    ("Tasks",      "Broken into\nactionable units"),
    ("GitHub Issues", "Each task = one\nissue, parallelizable"),
    ("Implementation","Sub-agent + worktree\nper issue, merge to main"),
]
sw_ = Inches(2.3); sh_ = Inches(1.6); gap = Inches(0.1)
total_w = sw_ * 5 + gap * 4
x_start = (SW - total_w) / 2
y = Inches(2.0)
for i, (t, b) in enumerate(stages):
    x = x_start + (sw_ + gap) * i
    fill = NAVY if i % 2 == 0 else NAVY_2
    add_round_rect(s, x, y, sw_, sh_, fill, line=fill, radius=0.1)
    add_text(s, x, y + Inches(0.25), sw_, Inches(0.5),
             t, size=15, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.8), sw_, Inches(0.7),
             b, size=10.5, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        # arrow chevron
        ax = x + sw_ + Emu(2000)
        ay = y + sh_ / 2 - Inches(0.12)
        chev = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, gap - Emu(4000), Inches(0.24))
        chev.shadow.inherit = False
        chev.fill.solid(); chev.fill.fore_color.rgb = CYAN
        chev.line.fill.background()

# Below the flow: artifacts + rituals
left_x = Inches(0.6); right_x = Inches(7.0); col_w_ = Inches(6.0)
add_text(s, left_x, Inches(4.1), col_w_, Inches(0.4),
         "Artifacts in the repo", size=15, bold=True, color=NAVY)
art = [
    "/prds/<feature>.md  —  product brief",
    "/epics/<feature>/  —  phased plan + task list",
    "GitHub Issues  —  one per task, label by epic",
    "/.claude/  —  agent + command definitions",
    "CLAUDE.md  —  project conventions for agents",
]
add_bullets(s, left_x, Inches(4.55), col_w_, Inches(2.5),
            art, size=11.5, gap=4, line_spacing=1.2)

add_text(s, right_x, Inches(4.1), col_w_, Inches(0.4),
         "Rituals (lightweight)", size=15, bold=True, color=NAVY)
rit = [
    "Spec-first.  No code before a PRD + epic exist.",
    "Issue-first.  Every change traces back to an issue.",
    "Branch-per-issue.  Worktree branches keep mainline green.",
    "Status sync.  Issue checkboxes are the truth.",
    "Pre-merge review.  Human signs off; AI doesn't merge itself.",
]
add_bullets(s, right_x, Inches(4.55), col_w_, Inches(2.5),
            rit, size=11.5, gap=4, line_spacing=1.2)
add_footer_brand(s)

# ── 11. Onboarding a new developer ───────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "09  ·  Onboarding", "A new developer's first month",
                page_no=11, total=TOTAL)

add_text(s, Inches(0.6), Inches(1.55), Inches(12.2), Inches(0.4),
         "What a fresh joiner does — and what they need from us — to be productive on Signature Sprint.",
         size=12, italic=True, color=MUTED)

phases = [
    ("Day 1", "Get oriented", [
        "Clone repo · run python -m http.server 8000 · play through both themes.",
        "Read CLAUDE.md, README, and PRODUCT-DOC §1 (PRD).",
        "Open the manager portal; read one report end-to-end with the Report-Reference at hand.",
    ]),
    ("Week 1", "Make a real change", [
        "Pick a Good-First-Issue from GitHub Issues.",
        "Spin a worktree, run Claude Code in it, ship a small PR end-to-end.",
        "Pair on one report-reading session with an evaluator on a flagged candidate.",
    ]),
    ("Month 1", "Own a vertical slice", [
        "Take an epic (e.g., new competency, persistence layer, locale pack).",
        "Author the PRD, decompose into issues, ship issue-by-issue.",
        "Present at the next internal review; update CLAUDE.md with anything you learned.",
    ]),
]
y = Inches(2.05)
for label, sub, items in phases:
    # left badge
    add_round_rect(s, Inches(0.6), y, Inches(2.2), Inches(1.55), NAVY, line=NAVY, radius=0.12)
    add_text(s, Inches(0.6), y + Inches(0.25), Inches(2.2), Inches(0.5),
             label, size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), y + Inches(0.8), Inches(2.2), Inches(0.5),
             sub, size=12, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    # right card
    add_round_rect(s, Inches(3.0), y, Inches(9.8), Inches(1.55), WHITE, line=RULE, radius=0.06)
    add_bullets(s, Inches(3.2), y + Inches(0.15), Inches(9.5), Inches(1.4),
                items, size=11.5, gap=2, line_spacing=1.2)
    y += Inches(1.7)

# ── 12. Tooling stack ────────────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "10  ·  Stack", "What CCPM looks like, concretely, for us",
                page_no=12, total=TOTAL)

cols = [
    ("Repo & code",
     [
        "GitHub  —  source of truth, Issues + PRs",
        "Branches via git worktree per issue",
        "CLAUDE.md  —  conventions, scoring, naming",
     ],
     CYAN),
    ("Claude Code",
     [
        "Per-issue sub-agent in worktree",
        "Custom slash-commands  —  /prd, /epic, /issue, /sync",
        "Hooks  —  pre-commit format, post-task issue update",
     ],
     CYAN_2),
    ("Docs & specs",
     [
        "PRODUCT-DOC.md  —  versioned PRD",
        "USER-GUIDE / REPORT-REFERENCE  —  user-facing",
        "/prds & /epics  —  feature-level briefs",
     ],
     NAVY_2),
    ("Quality gates",
     [
        "Manual review on every PR",
        "Print-rendered manager report check",
        "Smoke test in both themes before merge",
     ],
     NAVY),
]
cw = Inches(3.0); ch = Inches(4.5); gap = Inches(0.15)
x0 = Inches(0.6); y0 = Inches(1.7)
for i, (title, items, accent) in enumerate(cols):
    x = x0 + (cw + gap) * i
    add_round_rect(s, x, y0, cw, ch, WHITE, line=RULE, radius=0.06)
    add_rect(s, x, y0, cw, Inches(0.5), accent)
    add_text(s, x, y0 + Inches(0.05), cw, Inches(0.4),
             title, size=14, bold=True, color=WHITE if accent == NAVY else NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.2), y0 + Inches(0.7), cw - Inches(0.35), ch - Inches(0.85),
                items, size=11, gap=6, line_spacing=1.25)

# Closing strip
add_round_rect(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.7), ACCENT, line=CYAN, radius=0.18)
add_runs(s, Inches(0.85), Inches(6.5), Inches(11.8), Inches(0.5),
         [("Adoption cost.  ", {"size": 12, "bold": True, "color": NAVY}),
          ("≈ 2–3 days of one-time setup (CLAUDE.md, slash-commands, agents). After that, the cost is ~5 minutes per new issue.",
           {"size": 12, "color": INK})],
         anchor=MSO_ANCHOR.MIDDLE)

# ── 13. Risks & guardrails ───────────────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "11  ·  Risks", "Risks and the guardrails that contain them",
                page_no=13, total=TOTAL)

rows = [
    ["Risk",                                "Why it matters",                                                  "Guardrail"],
    ["AI confidently wrong",                "Plausible-looking code or scoring rules that are subtly off",     "Mandatory human review; tests on the scoring engine"],
    ["Context drift over long sessions",    "Decisions made early get \"forgotten\" later",                    "Capture in PRD / CLAUDE.md, not chat history"],
    ["Issue sprawl",                        "Many small issues with no epic = chaos",                          "Every issue links to an epic; no orphan issues"],
    ["Behavioral validity",                 "Game UX should not bend assessment science",                       "I/O-psych review on every new competency before launch"],
    ["Privacy / data handling",             "Real candidate data once we add a backend",                       "DPA + retention policy before persistence layer ships"],
    ["Single-engineer concentration",       "Bus factor of one",                                                "Onboard a second dev via the Day-1 / Week-1 / Month-1 path"],
]
add_clean_table(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(5.0), rows,
                col_widths=[Inches(3.0), Inches(4.5), Inches(4.7)],
                alt_fill=SOFT, body_size=11, header_size=11.5)

add_round_rect(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.5), NAVY, line=NAVY, radius=0.2)
add_runs(s, Inches(0.85), Inches(6.8), Inches(11.8), Inches(0.4),
         [("Posture.  ", {"size": 12, "bold": True, "color": CYAN}),
          ("AI accelerates; review decides. Speed-up doesn't change who is accountable for what ships.",
           {"size": 12, "color": WHITE})],
         anchor=MSO_ANCHOR.MIDDLE)

# ── 14. Roadmap (Now / Next / Later) ─────────────────────────────────────────
s = add_slide()
add_page_chrome(s, "12  ·  Roadmap", "Now · Next · Later",
                page_no=14, total=TOTAL)

cols = [
    ("Now",   "Next 2–3 weeks",
     ["Set up CCPM scaffolding (CLAUDE.md, /prds, /epics, slash-commands).",
      "Migrate the v2.7 backlog into GitHub Issues, grouped by epic.",
      "Run one full epic end-to-end as the reference example."],
     CYAN),
    ("Next",  "Quarter ahead",
     ["Ship persistence layer (Postgres + minimal REST).",
      "Add 1–2 new competencies (Collaboration · Customer Focus).",
      "Onboard a second engineer using the Day-1 / Week-1 / Month-1 path."],
     CYAN_2),
    ("Later", "Beyond this quarter",
     ["Multi-evaluator review with comments, shortlists, ATS export.",
      "Localized scenario packs (HI / regional).",
      "Cohort analytics + benchmarking dashboard."],
     NAVY),
]
cw = Inches(4.0); ch = Inches(5.0); gap = Inches(0.15)
x0 = Inches(0.6); y0 = Inches(1.6)
for i, (label, sub, items, accent) in enumerate(cols):
    x = x0 + (cw + gap) * i
    add_round_rect(s, x, y0, cw, ch, WHITE, line=RULE, radius=0.06)
    add_rect(s, x, y0, cw, Inches(0.7), accent)
    add_text(s, x + Inches(0.25), y0 + Inches(0.08), cw - Inches(0.4), Inches(0.35),
             label.upper(), size=11, bold=True, color=WHITE if accent == NAVY else NAVY)
    add_text(s, x + Inches(0.25), y0 + Inches(0.4), cw - Inches(0.4), Inches(0.3),
             sub, size=10.5, italic=True, color=WHITE if accent == NAVY else NAVY)
    add_bullets(s, x + Inches(0.25), y0 + Inches(0.95), cw - Inches(0.5), ch - Inches(1.2),
                items, size=11.5, gap=8, line_spacing=1.3)
add_footer_brand(s)

# ── 15. Closing ──────────────────────────────────────────────────────────────
s = add_slide()
add_rect(s, 0, 0, SW, SH, WHITE)
add_rect(s, 0, 0, SW, Inches(2.4), NAVY)
add_rect(s, 0, Inches(2.4), SW, Inches(0.05), CYAN)

# Wordmark big
tb = s.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(5), Inches(1.2))
p = tb.text_frame.paragraphs[0]
r1 = p.add_run(); r1.text = "hy"; r1.font.name = DISPLAY; r1.font.size = Pt(56); r1.font.bold = True; r1.font.color.rgb = WHITE
r2 = p.add_run(); r2.text = "reo"; r2.font.name = DISPLAY; r2.font.size = Pt(56); r2.font.bold = True; r2.font.color.rgb = CYAN

add_text(s, Inches(0.6), Inches(1.85), Inches(10), Inches(0.4),
         "Hyreo Labs  ·  Building with Claude.ai",
         size=13, italic=True, color=RGBColor(0xC7, 0xD7, 0xEA))

# Big closing line
add_text(s, Inches(0.6), Inches(3.0), Inches(12.2), Inches(1.0),
         "Speed comes from AI.\nDirection comes from us.",
         size=36, bold=True, color=NAVY)

# Three takeaways
t_y = Inches(4.7); cw = Inches(4.0); gap = Inches(0.15)
takeaways = [
    ("≈ 4–5×",       "Effort compression on this build",      CYAN),
    ("CCPM",         "The method we adopt going forward",     CYAN_2),
    ("Day 1 / W1 / M1","Onboarding path for new engineers",   NAVY_2),
]
for i, (big, sub, accent) in enumerate(takeaways):
    x = Inches(0.6) + i * (cw + gap)
    add_round_rect(s, x, t_y, cw, Inches(1.6), SOFT, line=RULE, radius=0.06)
    add_rect(s, x, t_y, Inches(0.12), Inches(1.6), accent)
    add_text(s, x + Inches(0.25), t_y + Inches(0.18), cw - Inches(0.4), Inches(0.7),
             big, size=26, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), t_y + Inches(0.95), cw - Inches(0.4), Inches(0.5),
             sub, size=12, color=INK)

add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "Discussion · questions · objections welcome.",
         size=14, italic=True, color=MUTED)
add_text(s, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
         "Hyreo Labs  ·  Internal  ·  Confidential",
         size=9, italic=True, color=MUTED)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"d:\BAT\Signature-Sprint-Building-with-Claude.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
